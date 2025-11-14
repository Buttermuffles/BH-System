#!/usr/bin/env python3
"""Lightweight local OCR fallback using Tesseract via pytesseract.

This script processes a file or directory and writes plain-text outputs
to the `ocr_output/` directory. It requires the Tesseract engine installed
on the system and the Python packages `pytesseract` and `Pillow`.

Supports: PDF, Images, Word (.doc, .docx), Excel (.xls, .xlsx), PowerPoint (.ppt, .pptx),
Text files (.txt), CSV, RTF, and OpenDocument formats.
"""
import argparse
import sys
import os
from pathlib import Path
import logging
from datetime import datetime, timezone
from dotenv import load_dotenv
import csv

try:
    from PIL import Image
    import pytesseract
except Exception as e:
    print("Missing dependencies. Install with: python -m pip install pytesseract pillow python-dotenv")
    raise

# Load environment variables from .env file (parent directory)
load_dotenv(Path(__file__).parent.parent / ".env")

# Set pytesseract path if TESSERACT_CMD is defined
# Try hardcoded path first as fallback
tesseract_cmd = os.getenv("TESSERACT_CMD") or r"C:\Program Files\Tesseract-OCR\tesseract.exe"
if tesseract_cmd and os.path.exists(tesseract_cmd):
    pytesseract.pytesseract.pytesseract_cmd = tesseract_cmd
    # Also set TESSDATA_PREFIX for language data files
    tessdata_dir = os.path.dirname(tesseract_cmd) + r"\tessdata"
    if os.path.exists(tessdata_dir):
        os.environ["TESSDATA_PREFIX"] = tessdata_dir


def setup_logging():
    ts = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
    logfile = Path("ocr_logs") / f"local_ocr_{ts}.log"
    logfile.parent.mkdir(parents=True, exist_ok=True)
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s %(levelname)s %(message)s",
        handlers=[
            logging.StreamHandler(sys.stdout),
            logging.FileHandler(logfile, encoding="utf-8"),
        ],
    )
    return logfile


def ocr_image(img_path: Path):
    logging.info(f"OCR image {img_path}")
    img = Image.open(img_path).convert("RGB")
    text = pytesseract.image_to_string(img)
    return text


def extract_text_from_docx(path: Path):
    """Extract text from .docx files."""
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [para.text for para in doc.paragraphs]
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    paragraphs.append(cell.text)
        return "\n".join(paragraphs)
    except ImportError:
        logging.warning("python-docx not installed. Install with: pip install python-docx")
        return None
    except Exception as e:
        logging.error(f"Failed to extract text from DOCX: {e}")
        return None


def extract_text_from_xlsx(path: Path):
    """Extract text from .xlsx files."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), data_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            texts.append(f"\n=== Sheet: {sheet_name} ===\n")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    texts.append(row_text)
        return "\n".join(texts)
    except ImportError:
        logging.warning("openpyxl not installed. Install with: pip install openpyxl")
        return None
    except Exception as e:
        logging.error(f"Failed to extract text from XLSX: {e}")
        return None


def extract_text_from_xls(path: Path):
    """Extract text from .xls files (old Excel format)."""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        texts = []
        for sheet_name in wb.sheet_names():
            sheet = wb.sheet_by_name(sheet_name)
            texts.append(f"\n=== Sheet: {sheet_name} ===\n")
            for row_idx in range(sheet.nrows):
                row = sheet.row(row_idx)
                row_text = "\t".join(str(cell.value) for cell in row)
                if row_text.strip():
                    texts.append(row_text)
        return "\n".join(texts)
    except ImportError:
        logging.warning("xlrd not installed. Install with: pip install xlrd")
        return None
    except Exception as e:
        logging.error(f"Failed to extract text from XLS: {e}")
        return None


def extract_text_from_pptx(path: Path):
    """Extract text from .pptx files."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts.append(f"\n=== Slide {slide_num} ===\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
    except ImportError:
        logging.warning("python-pptx not installed. Install with: pip install python-pptx")
        return None
    except Exception as e:
        logging.error(f"Failed to extract text from PPTX: {e}")
        return None


def extract_text_from_txt(path: Path):
    """Extract text from plain text files."""
    try:
        # Try UTF-8 first, then fallback to other encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                return path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        logging.warning(f"Could not decode {path} with any encoding")
        return None
    except Exception as e:
        logging.error(f"Failed to read text file: {e}")
        return None


def extract_text_from_csv(path: Path):
    """Extract text from CSV files."""
    try:
        texts = []
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                with open(path, 'r', encoding=encoding, newline='') as f:
                    reader = csv.reader(f)
                    for row in reader:
                        texts.append("\t".join(row))
                return "\n".join(texts)
            except UnicodeDecodeError:
                continue
        logging.warning(f"Could not decode CSV {path} with any encoding")
        return None
    except Exception as e:
        logging.error(f"Failed to read CSV file: {e}")
        return None


def extract_text_from_rtf(path: Path):
    """Extract text from RTF files."""
    try:
        from striprtf.striprtf import rtf_to_text
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            rtf_content = f.read()
        return rtf_to_text(rtf_content)
    except ImportError:
        logging.warning("striprtf not installed. Install with: pip install striprtf")
        return None
    except Exception as e:
        logging.error(f"Failed to extract text from RTF: {e}")
        return None


def process_file(path: Path, out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)
    suffix = path.suffix.lower()
    text = None
    
    # Handle different file types
    if suffix == ".pdf":
        # For PDFs, try to use pdf2image if available; otherwise, skip.
        try:
            from pdf2image import convert_from_path
            
            # Try to get poppler path from environment if configured
            poppler_path = os.getenv("POPPLER_PATH")
            if poppler_path:
                # Normalize path (handle forward slashes on Windows)
                poppler_path = os.path.normpath(poppler_path)
                if os.path.exists(poppler_path):
                    # Pass poppler path to convert_from_path
                    pages = convert_from_path(str(path), poppler_path=poppler_path)
                else:
                    # Try without explicit path (assumes poppler is in PATH)
                    pages = convert_from_path(str(path))
            else:
                # Try without explicit path (assumes poppler is in PATH)
                pages = convert_from_path(str(path))
            texts = []
            for i, p in enumerate(pages, start=1):
                t = pytesseract.image_to_string(p)
                texts.append(t)
            text = "\n\n".join(texts)
        except Exception as e:
            error_msg = str(e)
            if "poppler" in error_msg.lower() or "Unable to get page count" in error_msg:
                logging.error(f"PDF processing failed: Poppler is required for PDF processing.")
                logging.error(f"Install poppler: Windows: winget install poppler")
                logging.error(f"Or set POPPLER_PATH in .env to poppler bin directory")
            else:
                logging.error(f"pdf2image not available or failed: {e}")
            return False
    elif suffix == ".docx":
        text = extract_text_from_docx(path)
        if text is None:
            logging.warning(f"Could not extract text from DOCX, trying OCR fallback")
            # Fallback: Convert to PDF/image and OCR (would need additional tools)
            return False
    elif suffix == ".doc":
        # Old .doc format - try to use python-docx (might not work)
        # Better to convert to .docx first or use LibreOffice
        logging.warning(".doc files require conversion. Install LibreOffice or convert to .docx")
        return False
    elif suffix == ".xlsx":
        text = extract_text_from_xlsx(path)
        if text is None:
            return False
    elif suffix == ".xls":
        text = extract_text_from_xls(path)
        if text is None:
            return False
    elif suffix == ".pptx":
        text = extract_text_from_pptx(path)
        if text is None:
            return False
    elif suffix == ".ppt":
        logging.warning(".ppt files require conversion. Install LibreOffice or convert to .pptx")
        return False
    elif suffix == ".txt":
        text = extract_text_from_txt(path)
        if text is None:
            return False
    elif suffix == ".csv":
        text = extract_text_from_csv(path)
        if text is None:
            return False
    elif suffix == ".rtf":
        text = extract_text_from_rtf(path)
        if text is None:
            return False
    elif suffix in {".odt", ".ods", ".odp"}:
        # OpenDocument formats - would need odfpy or LibreOffice conversion
        logging.warning(f"OpenDocument format {suffix} requires conversion. Install LibreOffice or convert to standard formats")
        return False
    else:
        # Assume it's an image and try OCR
        try:
            text = ocr_image(path)
        except Exception as e:
            logging.error(f"Failed to process {path} as image: {e}")
            return False
    
    if text is None or text.strip() == "":
        logging.warning(f"No text extracted from {path}")
        return False

    out_file = out_dir / f"{path.stem}.txt"
    out_file.write_text(text, encoding="utf-8")
    logging.info(f"Wrote {out_file}")
    return True


def main():
    parser = argparse.ArgumentParser(description="Local OCR using pytesseract")
    parser.add_argument("--input", "-i", required=True, help="File or directory to process")
    parser.add_argument("--output", "-o", default="ocr_output", help="Output directory")
    args = parser.parse_args()

    log = setup_logging()
    logging.info("Starting local OCR")

    inp = Path(args.input)
    out_dir = Path(args.output)

    paths = []
    if inp.is_file():
        paths = [inp]
    elif inp.is_dir():
        # Supported file extensions
        exts = {
            # Images
            ".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif", ".webp",
            # Documents
            ".pdf", ".doc", ".docx", ".txt", ".rtf",
            # Spreadsheets
            ".xls", ".xlsx", ".csv",
            # Presentations
            ".ppt", ".pptx",
            # OpenDocument
            ".odt", ".ods", ".odp"
        }
        paths = [p for p in sorted(inp.rglob("*")) if p.is_file() and p.suffix.lower() in exts]
    else:
        logging.error(f"Input path not found: {inp}")
        sys.exit(1)

    summary = {"processed": 0, "failed": 0}
    for p in paths:
        try:
            ok = process_file(p, out_dir)
            if ok:
                summary["processed"] += 1
            else:
                summary["failed"] += 1
        except Exception as e:
            logging.exception(f"Failed {p}: {e}")
            summary["failed"] += 1

    logging.info(f"Done. Processed: {summary['processed']}, Failed: {summary['failed']}")


if __name__ == "__main__":
    main()
